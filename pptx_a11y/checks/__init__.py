from pptx.enum.shapes import MSO_SHAPE_TYPE


def iter_shapes(prs):
    """Yield (slide_index, shape) for every shape, descending into groups."""
    for i, slide in enumerate(prs.slides):
        yield from _walk(i, slide.shapes)


def _walk(slide_index, shapes):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _walk(slide_index, shape.shapes)
        else:
            yield slide_index, shape


ALL_CHECKS = []  # populated by register() below


def register(fn):
    ALL_CHECKS.append(fn)
    return fn


def load_all():
    """Import every check module so each @register call runs. Idempotent."""
    from importlib import import_module
    for name in (
        "alt_text", "slide_titles", "metadata", "contrast",
        "font_size", "link_text", "tables", "reading_order", "media_captions",
        "motion", "sensory", "use_of_color",
    ):
        import_module(f"{__name__}.{name}")
    return ALL_CHECKS
