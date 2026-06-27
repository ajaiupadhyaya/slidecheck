from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx_a11y.checks import iter_shapes
from pptx_a11y.fixers import register
from pptx_a11y.imageutil import image_bytes_and_type
from pptx_a11y.models import Change
from pptx_a11y.refs import shape_ref


def _alt(shape) -> str:
    try:
        return shape._element._nvXxPr.cNvPr.get("descr") or ""
    except Exception:  # noqa: BLE001
        return ""


def _slide_context(prs, slide_index: int, max_chars: int = 300) -> str:
    """Concatenated text from the slide, so the model can write alt text that
    fits the surrounding content (falls back to 'slide N')."""
    try:
        slide = prs.slides[slide_index]
    except IndexError:
        return f"slide {slide_index + 1}"
    parts: list[str] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            t = shape.text_frame.text.strip()
            if t:
                parts.append(t)
    text = " ".join(parts).strip()
    return text[:max_chars] if text else f"slide {slide_index + 1}"


@register
def fix(prs, describer) -> list[Change]:
    changes = []
    for slide_index, shape in iter_shapes(prs):
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        if _alt(shape).strip():
            continue
        payload = image_bytes_and_type(shape)
        if payload is None:
            continue
        blob, media_type = payload
        desc = describer.describe(blob, media_type, _slide_context(prs, slide_index))
        if not desc:
            continue
        shape._element._nvXxPr.cNvPr.set("descr", desc)
        changes.append(
            Change(
                fixer_id="alt_text",
                slide_index=slide_index,
                shape_ref=shape_ref(slide_index, shape),
                description=f'Added alt text: "{desc}"',
                machine_generated=True,
            )
        )
    return changes
