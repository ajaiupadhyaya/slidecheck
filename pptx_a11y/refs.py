from __future__ import annotations


def shape_ref(slide_index: int, shape) -> str:
    """Stable, human-readable reference for a shape within a slide."""
    return f"slide{slide_index}:shape{shape.shape_id}"


def shape_target(slide_index: int, shape) -> dict:
    """Return a stable machine-readable target dict for a shape.

    Used by fixers to locate the shape that needs remediation.
    """
    return {"slide": slide_index, "shape_id": shape.shape_id}


def run_target(slide_index: int, shape, para_index: int, run_index: int) -> dict:
    """Return a stable machine-readable target dict for a run within a shape.

    Para and run are integer indices into shape.text_frame.paragraphs[p].runs[r].
    """
    return {
        "slide": slide_index,
        "shape_id": shape.shape_id,
        "para": para_index,
        "run": run_index,
    }


def _find_shape_in_group(shapes, shape_id: int):
    """Recursively search a shape collection (including groups) for shape_id.

    Returns the matching shape or None.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape in shapes:
        if shape.shape_id == shape_id:
            return shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            found = _find_shape_in_group(shape.shapes, shape_id)
            if found is not None:
                return found
    return None


def resolve_target(prs, target: dict):
    """Resolve a *target* dict back to the live python-pptx object.

    Rules (checked in order):
    - ``{"scope": "document", ...}`` → returns *prs*
    - ``{"slide": i, "scope": "slide_title"}`` → returns *prs* (no shape to pin)
    - ``{"slide": i, "shape_id": id}`` → returns the matching shape (group-aware)
    - ``{"slide": i, "shape_id": id, "para": p, "run": r}`` → returns the run
    - Anything unresolvable → ``None``
    """
    # Document-scope targets (metadata, language, slide_title)
    if target.get("scope") in ("document", "slide_title"):
        return prs

    slide_idx = target.get("slide")
    shape_id = target.get("shape_id")

    if slide_idx is None or shape_id is None:
        return None

    try:
        slide = prs.slides[slide_idx]
    except IndexError:
        return None

    shape = _find_shape_in_group(slide.shapes, shape_id)
    if shape is None:
        return None

    para_idx = target.get("para")
    run_idx = target.get("run")

    if para_idx is None or run_idx is None:
        return shape

    # Run-level resolution
    try:
        run = shape.text_frame.paragraphs[para_idx].runs[run_idx]
        return run
    except (IndexError, AttributeError):
        return None
