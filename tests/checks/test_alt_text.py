import io

from lxml import etree
from pptx import Presentation
from pptx.util import Inches

from pptx_a11y.checks.alt_text import check
from tests.fixtures.build import _RED_PNG, clean_deck, deck_with_issues

_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_DEC_NS = "http://schemas.microsoft.com/office/drawing/2017/decorative"


def test_flags_picture_without_alt_text(tmp_path):
    prs = Presentation(deck_with_issues(str(tmp_path / "x.pptx")))
    findings = check(prs)
    assert any(f.check_id == "alt_text" and f.slide_index == 0 for f in findings)
    assert findings[0].severity.value == "error"
    # metadata assertions
    hit = next(f for f in findings if f.check_id == "alt_text")
    assert hit.fix_action == "set_alt_text"
    assert hit.fixable is True
    assert "shape_id" in hit.target
    assert hit.sc_refs == ["1.1.1"]
    assert hit.category == "images"


def test_clean_deck_has_no_alt_text_findings(tmp_path):
    prs = Presentation(clean_deck(str(tmp_path / "ok.pptx")))
    assert check(prs) == []


def test_chart_without_alt_text_gets_manual_hint(tmp_path):
    from pptx.util import Inches
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    data = CategoryChartData()
    data.categories = ["a", "b"]
    data.add_series("s", (1, 2))
    s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(4), Inches(3), data)
    findings = check(prs)
    chart_findings = [f for f in findings if f.check_id == "alt_text"]
    assert len(chart_findings) == 1
    assert "Chart" in chart_findings[0].message
    assert "manually" in chart_findings[0].suggestion.lower()


def test_native_decorative_image_not_flagged(tmp_path):
    """An image marked decorative via PowerPoint's native mechanism (an
    <adec:decorative> extension) needs no alt text and must not be flagged,
    even though its descr is empty."""
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    pic = s.shapes.add_picture(io.BytesIO(_RED_PNG), Inches(1), Inches(1), Inches(1), Inches(1))
    cNvPr = pic._element._nvXxPr.cNvPr
    cNvPr.set("descr", "")  # no alt text -> would normally be an ERROR
    ext_lst = etree.SubElement(cNvPr, f"{{{_A_NS}}}extLst")
    ext = etree.SubElement(ext_lst, f"{{{_A_NS}}}ext")
    ext.set("uri", "{C183D7F6-B498-43B3-948B-1728B52AA6E4}")
    dec = etree.SubElement(ext, f"{{{_DEC_NS}}}decorative")
    dec.set("val", "1")

    assert not any(f.check_id == "alt_text" for f in check(prs))
