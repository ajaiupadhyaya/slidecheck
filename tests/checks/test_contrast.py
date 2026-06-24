from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx_a11y.checks.contrast import check


def _deck_low_contrast(path):
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tf = s.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text_frame
    run = tf.paragraphs[0].add_run()
    run.text = "low contrast"
    run.font.size = Pt(24)
    run.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)  # light grey on default white
    prs.save(path)
    return path


def test_flags_low_contrast_with_ratio_and_suggestion(tmp_path):
    prs = Presentation(_deck_low_contrast(str(tmp_path / "c.pptx")))
    findings = check(prs)
    hits = [f for f in findings if f.check_id == "contrast"]
    assert hits
    assert "ratio" in hits[0].message.lower()
    assert hits[0].suggestion is not None


def test_no_false_positive_on_explicit_dark_box(tmp_path):
    """White text on a dark-filled shape: ~15:1 ratio — must NOT produce a contrast ERROR."""
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    box = s.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    # Set the shape fill to a dark colour so _background_rgb resolves it
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x20, 0x20, 0x20)
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "white on dark"
    run.font.size = Pt(24)
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    prs.save(str(tmp_path / "dark_box.pptx"))

    prs2 = Presentation(str(tmp_path / "dark_box.pptx"))
    findings = check(prs2)
    error_hits = [f for f in findings if f.check_id == "contrast" and f.severity.name == "ERROR"]
    assert not error_hits, f"False positive: got ERROR findings {error_hits}"


def test_flags_low_contrast_on_explicit_light_box(tmp_path):
    """Light-gray text on a white-filled shape: low contrast — must produce an ERROR
    and the message must NOT say 'assuming' since the background was explicitly resolved."""
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    box = s.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    # Explicit white fill on the shape
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "light gray on white"
    run.font.size = Pt(24)
    run.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
    prs.save(str(tmp_path / "light_box.pptx"))

    prs2 = Presentation(str(tmp_path / "light_box.pptx"))
    findings = check(prs2)
    error_hits = [f for f in findings if f.check_id == "contrast" and f.severity.name == "ERROR"]
    assert error_hits, "Expected a contrast ERROR for low-contrast light text on white box"
    assert "assuming" not in error_hits[0].message.lower(), (
        f"Message should not say 'assuming' when bg was resolved; got: {error_hits[0].message!r}"
    )
