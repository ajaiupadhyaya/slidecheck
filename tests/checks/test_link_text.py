from pptx import Presentation
from pptx_a11y.checks.link_text import check
from tests.fixtures.build import deck_with_issues


def test_flags_click_here(tmp_path):
    prs = Presentation(deck_with_issues(str(tmp_path / "x.pptx")))
    findings = check(prs)
    assert any(f.check_id == "link_text" for f in findings)
    # metadata assertions
    hit = next(f for f in findings if f.check_id == "link_text")
    assert hit.fix_action == "set_link_text"
    assert hit.fixable is True
    assert hit.current_value == "click here"
    assert "shape_id" in hit.target and "para" in hit.target and "run" in hit.target
    assert hit.sc_refs == ["2.4.4"]


def _deck_with_link(path, text):
    from pptx.util import Inches
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tf = s.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text_frame
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.hyperlink.address = "https://example.com/page"
    prs.save(path)
    return path


def test_flags_naked_domain_link_text(tmp_path):
    prs = Presentation(_deck_with_link(str(tmp_path / "d.pptx"), "example.com"))
    assert any(f.check_id == "link_text" for f in check(prs))


def test_descriptive_link_text_is_ok(tmp_path):
    prs = Presentation(_deck_with_link(str(tmp_path / "ok.pptx"), "Our research overview"))
    assert not any(f.check_id == "link_text" for f in check(prs))
