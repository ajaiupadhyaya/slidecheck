from pptx import Presentation
from pptx_a11y.checks.metadata import check


def test_flags_missing_core_title(tmp_path):
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    p = str(tmp_path / "m.pptx")
    prs.save(p)
    findings = check(Presentation(p))
    assert any("title" in f.message.lower() for f in findings)
    # metadata assertions
    title_hit = next(f for f in findings if f.shape_ref == "doc:title")
    assert title_hit.fix_action == "set_doc_title"
    assert title_hit.fixable is True
    assert title_hit.target == {"scope": "document", "field": "title"}
    assert title_hit.sc_refs == ["2.4.2"]


def test_title_present_not_flagged_for_title(tmp_path):
    prs = Presentation()
    prs.core_properties.title = "Has Title"
    prs.slides.add_slide(prs.slide_layouts[6])
    p = str(tmp_path / "m.pptx")
    prs.save(p)
    findings = check(Presentation(p))
    assert not any(f.shape_ref == "doc:title" for f in findings)


def test_flags_missing_language(tmp_path):
    prs = Presentation()
    prs.core_properties.title = "Has Title"  # title set, language still blank
    prs.slides.add_slide(prs.slide_layouts[6])
    p = str(tmp_path / "m.pptx")
    prs.save(p)
    findings = check(Presentation(p))
    assert any(f.shape_ref == "doc:language" for f in findings)


def test_language_present_not_flagged(tmp_path):
    prs = Presentation()
    prs.core_properties.title = "Has Title"
    prs.core_properties.language = "en-US"
    prs.slides.add_slide(prs.slide_layouts[6])
    p = str(tmp_path / "m.pptx")
    prs.save(p)
    findings = check(Presentation(p))
    assert findings == []
