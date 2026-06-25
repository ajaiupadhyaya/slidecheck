"""Tests for the motion / auto-advance check (WCAG 2.2.2)."""
from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn

from pptx_a11y.checks.motion import check

_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


def _inject_transition(slide, adv_tm: str | None = None, adv_click: str | None = None):
    """Add a <p:transition> element to *slide*._element with the given attrs."""
    trans = etree.SubElement(slide._element, qn("p:transition"))
    if adv_tm is not None:
        trans.set("advTm", adv_tm)
    if adv_click is not None:
        trans.set("advClick", adv_click)
    return trans


# ---------------------------------------------------------------------------
# RED → GREEN tests
# ---------------------------------------------------------------------------

def test_flags_slide_with_adv_tm():
    """A slide with advTm set must produce a motion finding."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _inject_transition(slide, adv_tm="3000")
    findings = check(prs)
    assert len(findings) == 1
    f = findings[0]
    assert f.check_id == "motion"
    assert f.slide_index == 0
    assert f.sc_refs == ["2.2.2"]
    assert f.fixable is False
    assert f.fix_action is None
    assert f.target == {"slide": 0, "scope": "slide"}
    assert f.section508 is True
    assert f.category == "motion"


def test_flags_slide_with_adv_click_zero():
    """advClick="0" alone also signals automatic advance."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _inject_transition(slide, adv_click="0")
    findings = check(prs)
    assert any(f.check_id == "motion" for f in findings)


def test_no_finding_without_transition():
    """A slide with no <p:transition> must not produce a motion finding."""
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    assert check(prs) == []


def test_no_finding_with_click_only_transition():
    """advClick="1" (click-to-advance) must not be flagged."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _inject_transition(slide, adv_click="1")
    assert check(prs) == []


def test_multiple_slides_only_auto_advance_flagged():
    """Only slides with auto-advance are flagged; others are clean."""
    prs = Presentation()
    s0 = prs.slides.add_slide(prs.slide_layouts[6])
    _inject_transition(s0, adv_tm="5000")
    prs.slides.add_slide(prs.slide_layouts[6])  # slide 1, no transition
    findings = check(prs)
    assert len(findings) == 1
    assert findings[0].slide_index == 0


def test_wcag_version_and_severity():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _inject_transition(slide, adv_tm="2000")
    f = check(prs)[0]
    from pptx_a11y.models import Severity
    assert f.severity == Severity.WARNING
    assert f.wcag_version == "2.0"
