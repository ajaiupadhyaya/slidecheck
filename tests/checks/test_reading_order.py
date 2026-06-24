from pptx import Presentation
from pptx.util import Inches

from pptx_a11y.checks.reading_order import check


def test_title_first_is_not_flagged(tmp_path):
    """Title placeholder is first in the shape tree -> good reading order.

    Guards the identity-bug fix: the old `shapes[0] is not title` compared
    fresh wrappers and was always True, so this case was falsely flagged.
    """
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    s.shapes.title.text = "Title"
    s.shapes.add_textbox(Inches(1), Inches(3), Inches(4), Inches(1)).text_frame.text = "body"

    assert not any(f.check_id == "reading_order" for f in check(prs))


def test_title_not_first_is_flagged(tmp_path):
    """Title moved to the end of the shape tree -> flagged INFO."""
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Title"
    s.shapes.add_textbox(Inches(1), Inches(3), Inches(4), Inches(1)).text_frame.text = "body"

    sp_tree = s.shapes._spTree
    title_el = s.shapes.title._element
    sp_tree.remove(title_el)
    sp_tree.append(title_el)

    assert any(f.check_id == "reading_order" for f in check(prs))
