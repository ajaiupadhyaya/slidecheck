"""Tests for the sensory-characteristics check (WCAG 1.3.3)."""
from pptx import Presentation
from pptx.util import Inches

from pptx_a11y.checks.sensory import check


def _deck_with_text(text: str) -> Presentation:
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tf = s.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text_frame
    tf.paragraphs[0].add_run().text = text
    return prs


# ---------------------------------------------------------------------------
# Positive cases (should flag)
# ---------------------------------------------------------------------------

def test_flags_position_instruction():
    prs = _deck_with_text("Click the button on the right to continue.")
    findings = check(prs)
    assert len(findings) == 1
    f = findings[0]
    assert f.check_id == "sensory"
    assert f.sc_refs == ["1.3.3"]
    assert f.fixable is False
    assert f.category == "text"
    assert f.section508 is True


def test_flags_click_green_button():
    prs = _deck_with_text("Click the green button to submit.")
    assert any(f.check_id == "sensory" for f in check(prs))


def test_flags_the_red_one():
    prs = _deck_with_text("Choose the red one from the list.")
    assert any(f.check_id == "sensory" for f in check(prs))


def test_case_insensitive():
    prs = _deck_with_text("CLICK THE BLUE BUTTON TO PROCEED.")
    assert any(f.check_id == "sensory" for f in check(prs))


# ---------------------------------------------------------------------------
# Negative cases (must NOT flag)
# ---------------------------------------------------------------------------

def test_see_below_not_flagged():
    """Normal lecture phrase 'See below for details.' must NOT be flagged."""
    prs = _deck_with_text("See below for details.")
    assert check(prs) == []


def test_shown_above_not_flagged():
    """Normal lecture phrase 'As shown above, ...' must NOT be flagged."""
    prs = _deck_with_text("As shown above, the chart confirms the trend.")
    assert check(prs) == []


def test_clean_text_no_finding():
    prs = _deck_with_text("Click the Submit button labelled 'OK' to continue.")
    assert check(prs) == []


def test_normal_prose_no_finding():
    prs = _deck_with_text("The results are summarised in the table on page 5.")
    assert check(prs) == []


def test_empty_deck_no_finding():
    prs = Presentation()
    assert check(prs) == []


# ---------------------------------------------------------------------------
# Deduplication: multiple matching runs in same shape → one finding
# ---------------------------------------------------------------------------

def test_deduplicates_per_shape():
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tf = s.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2)).text_frame
    # Add two paragraphs each with a sensory-only run
    tf.paragraphs[0].add_run().text = "Click the box on the left."
    p2 = tf.add_paragraph()
    p2.add_run().text = "Click the green button to confirm."
    findings = check(prs)
    # Should only produce ONE finding for this shape (deduplication by shape)
    assert len(findings) == 1
    assert findings[0].check_id == "sensory"
