from pptx import Presentation

from pptx_a11y.checks.slide_titles import check, check_title_quality
from tests.fixtures.build import clean_deck, deck_with_issues


def test_flags_slide_without_title(tmp_path):
    prs = Presentation(deck_with_issues(str(tmp_path / "x.pptx")))
    findings = check(prs)
    assert any(f.slide_index == 0 and f.check_id == "slide_title" for f in findings)
    # metadata assertions
    hit = next(f for f in findings if f.check_id == "slide_title")
    assert hit.fix_action == "set_title"
    assert hit.fixable is True
    assert hit.sc_refs == ["2.4.2"]
    assert hit.target.get("scope") == "slide_title"


def test_clean_deck_titles_ok(tmp_path):
    prs = Presentation(clean_deck(str(tmp_path / "ok.pptx")))
    assert check(prs) == []


# ---------------------------------------------------------------------------
# title_quality: generic titles
# ---------------------------------------------------------------------------

def _deck_with_titles(*titles: str) -> Presentation:
    """Create a deck with one slide per title (uses layout 5 = title-only)."""
    prs = Presentation()
    for t in titles:
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = t
    return prs


def test_flags_generic_title_slide():
    prs = _deck_with_titles("Slide 1")
    findings = check_title_quality(prs)
    assert any(f.check_id == "title_quality" and "generic" in f.message for f in findings)


def test_flags_generic_title_untitled():
    prs = _deck_with_titles("Untitled")
    findings = check_title_quality(prs)
    assert any(f.check_id == "title_quality" and "generic" in f.message for f in findings)


def test_flags_generic_title_presentation():
    prs = _deck_with_titles("Presentation")
    findings = check_title_quality(prs)
    assert any(f.check_id == "title_quality" and "generic" in f.message for f in findings)


def test_flags_generic_title_case_insensitive():
    prs = _deck_with_titles("SLIDE 3")
    findings = check_title_quality(prs)
    assert any(f.check_id == "title_quality" and "generic" in f.message for f in findings)


def test_descriptive_title_not_flagged_as_generic():
    prs = _deck_with_titles("Introduction to Machine Learning")
    findings = check_title_quality(prs)
    assert not any("generic" in f.message for f in findings)


# ---------------------------------------------------------------------------
# title_quality: duplicate titles
# ---------------------------------------------------------------------------

def test_flags_duplicate_titles():
    prs = _deck_with_titles("Overview", "Overview")
    findings = check_title_quality(prs)
    dups = [f for f in findings if "duplicated" in f.message]
    # Both slides should be flagged
    assert len(dups) == 2
    assert {f.slide_index for f in dups} == {0, 1}


def test_flags_duplicate_titles_case_insensitive():
    prs = _deck_with_titles("overview", "Overview")
    findings = check_title_quality(prs)
    assert any("duplicated" in f.message for f in findings)


def test_unique_titles_not_flagged_as_duplicates():
    prs = _deck_with_titles("Introduction", "Methods", "Results")
    findings = check_title_quality(prs)
    assert not any("duplicated" in f.message for f in findings)


# ---------------------------------------------------------------------------
# title_quality: metadata assertions
# ---------------------------------------------------------------------------

def test_title_quality_metadata():
    prs = _deck_with_titles("Slide 1")
    f = check_title_quality(prs)[0]
    assert f.check_id == "title_quality"
    assert f.sc_refs == ["2.4.6"]
    assert f.wcag_version == "2.0"
    assert f.section508 is True
    assert f.fixable is True
    assert f.fix_action == "set_title"
    assert f.target == {"slide": 0, "scope": "slide_title"}
    assert f.current_value == "Slide 1"


# ---------------------------------------------------------------------------
# title_quality: weak-title heuristics (long / all-caps / numeric-only)
# ---------------------------------------------------------------------------

def test_flags_overly_long_title():
    prs = _deck_with_titles("A" * 85)
    findings = check_title_quality(prs)
    assert any("long" in f.message.lower() for f in findings)


def test_does_not_flag_normal_length_title_as_long():
    prs = _deck_with_titles("Introduction to Machine Learning")
    findings = check_title_quality(prs)
    assert not any("long" in f.message.lower() for f in findings)


def test_flags_all_caps_title_without_destructive_suggestion():
    # All-caps titles often contain acronyms ("NASA BUDGET 2024"); a naive
    # .title() would mangle them ("Nasa Budget 2024"). So we flag but do NOT
    # offer an auto-fillable suggestion — the human rewrites it.
    prs = _deck_with_titles("NASA BUDGET 2024")
    findings = check_title_quality(prs)
    hit = next((f for f in findings if "capital" in f.message.lower()), None)
    assert hit is not None
    assert hit.suggested_value is None


def test_does_not_flag_short_acronym_title_as_all_caps():
    # single-word acronyms (FAQ, API) are legitimate titles, not shouting.
    prs = _deck_with_titles("FAQ")
    findings = check_title_quality(prs)
    assert not any("capital" in f.message.lower() for f in findings)


def test_flags_numeric_only_title():
    prs = _deck_with_titles("2.3")
    findings = check_title_quality(prs)
    assert any("number" in f.message.lower() for f in findings)


def test_descriptive_mixedcase_title_has_no_quality_findings():
    prs = _deck_with_titles("Results and Discussion")
    assert check_title_quality(prs) == []


def test_empty_deck_no_title_quality_findings():
    prs = Presentation()
    assert check_title_quality(prs) == []


def test_slides_without_titles_not_double_counted():
    """Slides with missing titles must not appear in title_quality findings."""
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = ""  # empty title — handled by check(), not check_title_quality()
    findings = check_title_quality(prs)
    # empty title → skip in title_quality (it's handled by check())
    assert findings == []
