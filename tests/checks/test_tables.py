from pptx import Presentation
from pptx.util import Inches
from pptx_a11y.checks.tables import check


def _deck_with_table(path, header):
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    gf = s.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2))
    tbl = gf.table
    tbl.first_row = header
    prs.save(path)
    return path


def test_flags_table_without_header_row(tmp_path):
    prs = Presentation(_deck_with_table(str(tmp_path / "t.pptx"), header=False))
    findings = check(prs)
    assert any(f.check_id == "table" for f in findings)
    # metadata assertions
    hit = next(f for f in findings if f.check_id == "table")
    assert hit.fix_action == "set_table_header"
    assert hit.fixable is True
    assert "shape_id" in hit.target
    assert hit.sc_refs == ["1.3.1"]


def test_table_with_header_row_ok(tmp_path):
    prs = Presentation(_deck_with_table(str(tmp_path / "t.pptx"), header=True))
    findings = check(prs)
    assert not any("header" in f.message.lower() for f in findings)
