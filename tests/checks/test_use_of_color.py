"""Tests for the use-of-color / link-underline check (WCAG 1.4.1)."""
from pptx import Presentation
from pptx.util import Inches

from pptx_a11y.checks.use_of_color import check


def _deck_with_link(text: str, underline: bool | None) -> Presentation:
    """Create a deck with one hyperlinked run, optionally underlined."""
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tf = s.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text_frame
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.hyperlink.address = "https://example.com"
    if underline is not None:
        run.font.underline = underline
    return prs


# ---------------------------------------------------------------------------
# Positive cases (should flag)
# ---------------------------------------------------------------------------

def test_flags_link_with_underline_explicitly_false():
    """Hyperlinked run with underline explicitly set to False must be flagged."""
    prs = _deck_with_link("Download here", underline=False)
    findings = check(prs)
    assert len(findings) == 1
    f = findings[0]
    assert f.check_id == "use_of_color"
    assert f.sc_refs == ["1.4.1"]
    assert f.fixable is False
    assert f.category == "color"
    assert f.section508 is True
    assert f.wcag_version == "2.0"
    assert "underline removed" in f.message


# ---------------------------------------------------------------------------
# Negative cases (must NOT flag)
# ---------------------------------------------------------------------------

def test_no_finding_for_link_with_underline_none():
    """Hyperlinked run with underline=None (inherited/theme default) must NOT be flagged.

    Office themes underline links by default; we cannot resolve the theme here,
    so None must be treated as 'may be underlined' — not a violation.
    """
    prs = _deck_with_link("Our research page", underline=None)
    assert check(prs) == []


def test_no_finding_for_underlined_link():
    """Hyperlinked run WITH underline must not be flagged."""
    prs = _deck_with_link("Read the full report", underline=True)
    assert check(prs) == []


def test_no_finding_for_plain_text_run():
    """Non-hyperlinked run must not be flagged."""
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tf = s.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text_frame
    run = tf.paragraphs[0].add_run()
    run.text = "Just plain text, no link."
    assert check(prs) == []


def test_no_finding_empty_deck():
    prs = Presentation()
    assert check(prs) == []


# ---------------------------------------------------------------------------
# Target / metadata assertions
# ---------------------------------------------------------------------------

def test_target_has_shape_id():
    prs = _deck_with_link("Visit our site", underline=False)
    f = check(prs)[0]
    assert "shape_id" in f.target
    assert f.target.get("slide") == 0


def test_severity_is_warning():
    from pptx_a11y.models import Severity
    prs = _deck_with_link("Visit our site", underline=False)
    f = check(prs)[0]
    assert f.severity == Severity.WARNING
