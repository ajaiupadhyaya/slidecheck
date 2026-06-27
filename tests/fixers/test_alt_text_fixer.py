from pptx import Presentation
from pptx_a11y.fixers.alt_text import fix
from tests.fixtures.build import _add_picture, deck_with_issues


class _StubDescriber:
    def describe(self, image_bytes, media_type, context):
        return "A solid red square."


class _RecordingDescriber:
    def __init__(self):
        self.contexts = []

    def describe(self, image_bytes, media_type, context):
        self.contexts.append(context)
        return "alt"


def test_alt_text_fix_passes_slide_text_as_context(tmp_path):
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Photosynthesis"
    _add_picture(s, None)
    rec = _RecordingDescriber()
    fix(prs, rec)
    assert rec.contexts, "describer should be called for the un-described picture"
    # the slide's own text is handed to the model so alt text is relevant
    assert "Photosynthesis" in rec.contexts[0]


class _NullDescriber:
    def describe(self, image_bytes, media_type, context):
        return None


def _first_picture(prs):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return shape
    return None


def test_alt_text_fix_embeds_description(tmp_path):
    prs = Presentation(deck_with_issues(str(tmp_path / "x.pptx")))
    changes = fix(prs, _StubDescriber())
    assert any(c.fixer_id == "alt_text" and c.machine_generated for c in changes)
    pic = _first_picture(prs)
    assert pic._element._nvXxPr.cNvPr.get("descr") == "A solid red square."


def test_alt_text_fix_no_change_when_describer_returns_none(tmp_path):
    prs = Presentation(deck_with_issues(str(tmp_path / "x.pptx")))
    changes = fix(prs, _NullDescriber())
    assert changes == []
    pic = _first_picture(prs)
    assert not (pic._element._nvXxPr.cNvPr.get("descr") or "")
