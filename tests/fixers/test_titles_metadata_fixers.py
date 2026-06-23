from pptx import Presentation
from pptx_a11y.fixers.slide_titles import fix as fix_titles
from pptx_a11y.fixers.metadata import fix as fix_metadata
from tests.fixtures.build import deck_with_issues


class _TitleDescriber:
    def describe(self, image_bytes, media_type, context):
        return "Overview"

    def suggest_text(self, prompt):
        return "Overview"


class _NullTitleDescriber:
    def describe(self, image_bytes, media_type, context):
        return None

    def suggest_text(self, prompt):
        return None


def test_title_fixer_fills_empty_title(tmp_path):
    prs = Presentation(deck_with_issues(str(tmp_path / "x.pptx")))
    changes = fix_titles(prs, _TitleDescriber())
    title_changes = [c for c in changes if c.fixer_id == "slide_title"]
    assert title_changes, "expected at least one slide_title change"
    filled_change = title_changes[0]
    assert filled_change.machine_generated is True
    assert prs.slides[0].shapes.title.text.strip() == "Overview"


def test_title_fixer_fallback_when_no_suggestion(tmp_path):
    prs = Presentation(deck_with_issues(str(tmp_path / "x.pptx")))
    changes = fix_titles(prs, _NullTitleDescriber())
    title_changes = [c for c in changes if c.fixer_id == "slide_title"]
    assert title_changes, "expected at least one slide_title change"
    fallback_change = title_changes[0]
    assert prs.slides[0].shapes.title.text.strip() == "Slide 1"
    assert fallback_change.machine_generated is False


def test_metadata_fixer_sets_title(tmp_path):
    prs = Presentation(deck_with_issues(str(tmp_path / "x.pptx")))
    prs.core_properties.title = ""
    changes = fix_metadata(prs, _TitleDescriber())
    assert any(c.fixer_id == "metadata" for c in changes)
    assert prs.core_properties.title.strip() != ""
