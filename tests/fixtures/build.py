import base64
import io
from pptx import Presentation
from pptx.util import Inches, Pt

# 1x1 red PNG
_RED_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNkYPhfDwAChwGA60e6kgAAAABJRU5ErkJggg=="
)


def _add_picture(slide, alt: str | None):
    pic = slide.shapes.add_picture(io.BytesIO(_RED_PNG), Inches(1), Inches(1), Inches(1), Inches(1))
    # Explicitly set descr to empty string if None, otherwise set to provided text
    pic._element._nvXxPr.cNvPr.set("descr", alt or "")
    return pic


def clean_deck(path: str) -> str:
    """A deck with no accessibility issues."""
    prs = Presentation()
    prs.core_properties.title = "Clean Deck"
    s = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    s.shapes.title.text = "Intro"
    _add_picture(s, "A red square")
    prs.save(path)
    return path


def deck_with_issues(path: str) -> str:
    """A deck planted with one of each issue the checks target."""
    prs = Presentation()
    # slide 0: missing title + picture with no alt text
    s0 = prs.slides.add_slide(prs.slide_layouts[5])
    s0.shapes.title.text = ""           # missing title
    _add_picture(s0, None)              # missing alt text
    # slide 1: tiny font + bad link text
    s1 = prs.slides.add_slide(prs.slide_layouts[5])
    s1.shapes.title.text = "Details"
    tb = s1.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1)).text_frame
    run = tb.paragraphs[0].add_run()
    run.text = "click here"
    run.font.size = Pt(10)             # too small
    # add a hyperlink so link_text check sees it
    run.hyperlink.address = "https://example.com"
    prs.save(path)
    return path
