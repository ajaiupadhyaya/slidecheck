"""Tests for pptx_a11y.analyze — TDD, RED then GREEN."""
from __future__ import annotations

import io

import pytest
from pptx import Presentation

from pptx_a11y.alt_text_ai import NullDescriber
from pptx_a11y.analyze import analyze, finding_to_dict, generate_suggestions, run_checks
from tests.fixtures.build import clean_deck, deck_with_issues


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

@pytest.fixture()
def issues_prs(tmp_path) -> Presentation:
    path = deck_with_issues(str(tmp_path / "issues.pptx"))
    return Presentation(path)


@pytest.fixture()
def clean_prs(tmp_path) -> Presentation:
    path = clean_deck(str(tmp_path / "clean.pptx"))
    return Presentation(path)


class _FakeDescriber:
    """Returns fixed strings so we can assert they flow into suggested_value."""

    FIXED_ALT = "A fake alt text from describer."
    FIXED_TEXT = "Fake Suggested Text"

    def describe(self, image_bytes: bytes, media_type: str, context: str) -> str | None:
        return self.FIXED_ALT

    def suggest_text(self, prompt: str) -> str | None:
        return self.FIXED_TEXT


# ---------------------------------------------------------------------------
# run_checks — does NOT mutate the deck
# ---------------------------------------------------------------------------

def test_run_checks_does_not_mutate_deck(tmp_path):
    """Checks must not change any user-visible content in the deck.

    python-pptx lazily creates some XML elements when background/fill
    properties are read (e.g. in the contrast check), so byte-level equality
    cannot be used.  Instead we verify that the semantically important fields
    the checks could plausibly modify are unchanged after run_checks().
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Pt

    path = deck_with_issues(str(tmp_path / "deck.pptx"))
    prs = Presentation(path)

    # Snapshot key content before checks.
    s0_title_before = prs.slides[0].shapes.title.text
    s1_title_before = prs.slides[1].shapes.title.text
    pics_before = {
        sh.shape_id: (sh._element._nvXxPr.cNvPr.get("descr") or "")
        for slide in prs.slides
        for sh in slide.shapes
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE
    }

    run_checks(prs)

    # Verify nothing changed.
    assert prs.slides[0].shapes.title.text == s0_title_before, "slide 0 title was mutated"
    assert prs.slides[1].shapes.title.text == s1_title_before, "slide 1 title was mutated"
    pics_after = {
        sh.shape_id: (sh._element._nvXxPr.cNvPr.get("descr") or "")
        for slide in prs.slides
        for sh in slide.shapes
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE
    }
    assert pics_before == pics_after, "picture alt texts were mutated"


def test_run_checks_returns_list_of_findings(issues_prs):
    findings = run_checks(issues_prs)
    assert isinstance(findings, list)
    assert len(findings) > 0


def test_run_checks_clean_deck_has_fewer_errors(clean_prs, issues_prs):
    clean_findings = run_checks(clean_prs)
    issue_findings = run_checks(issues_prs)
    clean_errors = sum(1 for f in clean_findings if f.severity.value == "error")
    issue_errors = sum(1 for f in issue_findings if f.severity.value == "error")
    assert issue_errors > clean_errors


# ---------------------------------------------------------------------------
# finding_to_dict — JSON-safe serialization
# ---------------------------------------------------------------------------

def test_finding_to_dict_has_required_keys(issues_prs):
    findings = run_checks(issues_prs)
    assert findings, "need at least one finding"
    d = finding_to_dict(findings[0])
    # Structural keys required by the task brief
    for key in ("id", "check_id", "severity", "slide_index", "message",
                 "sc_refs", "fix_action", "target", "fixable",
                 "suggested_value", "current_value", "shape_ref",
                 "auto_fixed", "suggestion", "wcag_version", "section508",
                 "category"):
        assert key in d, f"finding_to_dict missing key: {key!r}"


def test_finding_to_dict_severity_is_string(issues_prs):
    findings = run_checks(issues_prs)
    for f in findings:
        d = finding_to_dict(f)
        assert isinstance(d["severity"], str), "severity must be serialized to string"


def test_finding_to_dict_id_is_stable_string(issues_prs):
    findings = run_checks(issues_prs)
    f = findings[0]
    d = finding_to_dict(f)
    assert isinstance(d["id"], str)
    # id must contain check_id and slide_index
    assert f.check_id in d["id"]
    assert str(f.slide_index) in d["id"]


def test_finding_to_dict_sc_refs_is_list(issues_prs):
    findings = run_checks(issues_prs)
    for f in findings:
        d = finding_to_dict(f)
        assert isinstance(d["sc_refs"], list)


# ---------------------------------------------------------------------------
# generate_suggestions — deterministic paths (NullDescriber)
# ---------------------------------------------------------------------------

def test_bump_font_size_gets_deterministic_suggestion(issues_prs):
    findings = run_checks(issues_prs)
    font_findings = [f for f in findings if f.fix_action == "bump_font_size"]
    assert font_findings, "deck_with_issues must have a font_size finding"

    generate_suggestions(issues_prs, findings, NullDescriber())
    for f in font_findings:
        assert f.suggested_value == "18", (
            f"bump_font_size should set suggested_value='18', got {f.suggested_value!r}"
        )


def test_set_doc_language_gets_deterministic_suggestion(issues_prs):
    findings = run_checks(issues_prs)
    lang_findings = [f for f in findings if f.fix_action == "set_doc_language"]
    assert lang_findings, "deck_with_issues must have a set_doc_language finding"

    generate_suggestions(issues_prs, findings, NullDescriber())
    for f in lang_findings:
        assert f.suggested_value == "en-US", (
            f"set_doc_language should set 'en-US', got {f.suggested_value!r}"
        )


def test_set_doc_title_fallback_is_presentation(tmp_path):
    """When describer returns None (NullDescriber), fall back to 'Presentation'."""
    # Build a deck with no title so metadata check fires.
    prs = Presentation()
    prs.core_properties.title = ""
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = ""
    path = str(tmp_path / "notitle.pptx")
    prs.save(path)
    prs2 = Presentation(path)

    findings = run_checks(prs2)
    title_findings = [f for f in findings if f.fix_action == "set_doc_title"]
    assert title_findings

    generate_suggestions(prs2, findings, NullDescriber())
    for f in title_findings:
        assert f.suggested_value == "Presentation", (
            f"Expected fallback 'Presentation', got {f.suggested_value!r}"
        )


def test_set_alt_text_is_none_under_null_describer(issues_prs):
    findings = run_checks(issues_prs)
    alt_findings = [f for f in findings if f.fix_action == "set_alt_text"]
    assert alt_findings

    generate_suggestions(issues_prs, findings, NullDescriber())
    for f in alt_findings:
        assert f.suggested_value is None, (
            f"NullDescriber should yield None for set_alt_text, got {f.suggested_value!r}"
        )


def test_set_title_is_none_under_null_describer(issues_prs):
    findings = run_checks(issues_prs)
    title_findings = [f for f in findings if f.fix_action == "set_title"]
    assert title_findings

    generate_suggestions(issues_prs, findings, NullDescriber())
    for f in title_findings:
        assert f.suggested_value is None, (
            f"NullDescriber should yield None for set_title, got {f.suggested_value!r}"
        )


def test_set_link_text_is_none_under_null_describer(issues_prs):
    findings = run_checks(issues_prs)
    link_findings = [f for f in findings if f.fix_action == "set_link_text"]
    assert link_findings

    generate_suggestions(issues_prs, findings, NullDescriber())
    for f in link_findings:
        assert f.suggested_value is None


def test_apply_contrast_color_preserves_existing_suggested_value(issues_prs):
    """Contrast check already sets suggested_value; generate_suggestions must not overwrite it."""
    findings = run_checks(issues_prs)
    contrast_findings = [
        f for f in findings
        if f.fix_action == "apply_contrast_color" and f.suggested_value is not None
    ]
    if not contrast_findings:
        pytest.skip("no contrast findings with pre-set suggested_value in this deck")

    original_vals = {id(f): f.suggested_value for f in contrast_findings}
    generate_suggestions(issues_prs, findings, NullDescriber())
    for f in contrast_findings:
        assert f.suggested_value == original_vals[id(f)], (
            "generate_suggestions must not overwrite existing suggested_value"
        )


# ---------------------------------------------------------------------------
# generate_suggestions — fake describer flows through
# ---------------------------------------------------------------------------

def test_set_alt_text_uses_describer(issues_prs):
    findings = run_checks(issues_prs)
    alt_findings = [f for f in findings if f.fix_action == "set_alt_text"]
    assert alt_findings

    generate_suggestions(issues_prs, findings, _FakeDescriber())
    for f in alt_findings:
        assert f.suggested_value == _FakeDescriber.FIXED_ALT, (
            f"expected describer text, got {f.suggested_value!r}"
        )


def test_set_title_uses_describer(issues_prs):
    findings = run_checks(issues_prs)
    title_findings = [f for f in findings if f.fix_action == "set_title"]
    assert title_findings

    generate_suggestions(issues_prs, findings, _FakeDescriber())
    for f in title_findings:
        assert f.suggested_value == _FakeDescriber.FIXED_TEXT


def test_set_link_text_uses_describer(issues_prs):
    findings = run_checks(issues_prs)
    link_findings = [f for f in findings if f.fix_action == "set_link_text"]
    assert link_findings

    generate_suggestions(issues_prs, findings, _FakeDescriber())
    for f in link_findings:
        assert f.suggested_value == _FakeDescriber.FIXED_TEXT


def test_set_doc_title_uses_describer(issues_prs):
    findings = run_checks(issues_prs)
    doc_title_findings = [f for f in findings if f.fix_action == "set_doc_title"]
    assert doc_title_findings

    generate_suggestions(issues_prs, findings, _FakeDescriber())
    for f in doc_title_findings:
        assert f.suggested_value == _FakeDescriber.FIXED_TEXT


# ---------------------------------------------------------------------------
# analyze — top-level integration
# ---------------------------------------------------------------------------

def test_analyze_returns_correct_structure(issues_prs):
    result = analyze(issues_prs, NullDescriber())
    assert "findings" in result
    assert "score" in result
    assert "coverage" in result


def test_analyze_findings_are_dicts(issues_prs):
    result = analyze(issues_prs, NullDescriber())
    assert isinstance(result["findings"], list)
    assert len(result["findings"]) > 0
    for d in result["findings"]:
        assert isinstance(d, dict)


def test_analyze_finding_dicts_have_required_keys(issues_prs):
    result = analyze(issues_prs, NullDescriber())
    for d in result["findings"]:
        for key in ("id", "sc_refs", "fix_action", "target", "severity"):
            assert key in d, f"missing key {key!r} in finding dict"


def test_analyze_score_has_expected_keys(issues_prs):
    result = analyze(issues_prs, NullDescriber())
    score = result["score"]
    assert "score" in score
    assert "grade" in score
    assert "errors" in score
    assert "warnings" in score


def test_analyze_score_is_int_between_0_and_100(issues_prs):
    result = analyze(issues_prs, NullDescriber())
    s = result["score"]["score"]
    assert isinstance(s, int)
    assert 0 <= s <= 100


def test_analyze_coverage_is_list_of_dicts(issues_prs):
    result = analyze(issues_prs, NullDescriber())
    cov = result["coverage"]
    assert isinstance(cov, list)
    assert len(cov) > 0
    for row in cov:
        assert "sc" in row
        assert "status" in row


def test_analyze_deterministic_suggestions_set(issues_prs):
    """bump_font_size and set_doc_language should have suggested_value in final dicts."""
    result = analyze(issues_prs, NullDescriber())
    font_dicts = [d for d in result["findings"] if d.get("fix_action") == "bump_font_size"]
    lang_dicts = [d for d in result["findings"] if d.get("fix_action") == "set_doc_language"]

    assert font_dicts, "expect at least one bump_font_size finding"
    assert lang_dicts, "expect at least one set_doc_language finding"

    for d in font_dicts:
        assert d["suggested_value"] == "18"
    for d in lang_dicts:
        assert d["suggested_value"] == "en-US"


def test_analyze_ai_suggestions_none_under_null_describer(issues_prs):
    result = analyze(issues_prs, NullDescriber())
    for d in result["findings"]:
        if d.get("fix_action") in ("set_alt_text", "set_title", "set_link_text"):
            assert d["suggested_value"] is None, (
                f"AI suggestion should be None under NullDescriber for {d['fix_action']!r}"
            )


def test_analyze_clean_deck_scores_higher(clean_prs, issues_prs):
    clean_result = analyze(clean_prs, NullDescriber())
    issue_result = analyze(issues_prs, NullDescriber())
    assert clean_result["score"]["score"] >= issue_result["score"]["score"]
