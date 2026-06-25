"""Tests for pptx_a11y.appliers — deterministic fix-applier registry."""
from __future__ import annotations

import io
import base64

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from tests.fixtures.build import deck_with_issues, clean_deck

# ---------------------------------------------------------------------------
# helpers
# ---------------------------------------------------------------------------

# 1×1 red PNG (same as build.py)
_RED_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNkYPhfDwAChwGA60e6kgAAAABJRU5ErkJggg=="
)


def _prs_with_picture(tmp_path):
    """Deck with a picture on slide 0 (no alt text set)."""
    prs = Presentation(deck_with_issues(str(tmp_path / "d.pptx")))
    return prs


def _prs_with_run(tmp_path):
    """Deck_with_issues: slide 1 has a run with a hyperlink."""
    prs = Presentation(deck_with_issues(str(tmp_path / "d.pptx")))
    return prs


def _pic_and_target(prs):
    """Return (pic_shape, shape_target_dict) from slide 0."""
    slide0 = prs.slides[0]
    pic = next(s for s in slide0.shapes if s.shape_type == 13)  # PICTURE
    return pic, {"slide": 0, "shape_id": pic.shape_id}


def _run_and_target(prs):
    """Return (run, run_target_dict) from slide 1 (the hyperlink run).

    Specifically targets the textbox run with 'click here' and a hyperlink,
    not the title placeholder.
    """
    slide1 = prs.slides[1]
    # Find a shape whose first run has a hyperlink (the textbox, not the title)
    shape = next(
        s for s in slide1.shapes
        if s.has_text_frame
        and s.text_frame.paragraphs
        and s.text_frame.paragraphs[0].runs
        and s.text_frame.paragraphs[0].runs[0].hyperlink.address is not None
    )
    run = shape.text_frame.paragraphs[0].runs[0]
    target = {
        "slide": 1,
        "shape_id": shape.shape_id,
        "para": 0,
        "run": 0,
    }
    return run, target


def _prs_with_table(tmp_path):
    """Deck with a table (first_row=False) on slide 0."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Table slide"
    tbl = slide.shapes.add_table(3, 3, Inches(1), Inches(1), Inches(4), Inches(2)).table
    # python-pptx may default first_row to True; explicitly disable it
    tbl.first_row = False
    path = str(tmp_path / "tbl.pptx")
    prs.save(path)
    prs2 = Presentation(path)
    # Verify the fixture is set up correctly
    tbl2 = next(s for s in prs2.slides[0].shapes if s.has_table).table
    assert not tbl2.first_row, "fixture setup failed: first_row should be False"
    return prs2


def _table_shape_target(prs):
    slide0 = prs.slides[0]
    shape = next(s for s in slide0.shapes if s.has_table)
    return shape, {"slide": 0, "shape_id": shape.shape_id}


# ---------------------------------------------------------------------------
# set_alt_text
# ---------------------------------------------------------------------------

def test_set_alt_text_sets_descr(tmp_path):
    from pptx_a11y.appliers import set_alt_text

    prs = _prs_with_picture(tmp_path)
    pic, target = _pic_and_target(prs)
    ok = set_alt_text(prs, target, "A red square image")
    assert ok is True
    assert pic._element._nvXxPr.cNvPr.get("descr") == "A red square image"


def test_set_alt_text_returns_false_for_unknown_target(tmp_path):
    from pptx_a11y.appliers import set_alt_text

    prs = _prs_with_picture(tmp_path)
    ok = set_alt_text(prs, {"slide": 0, "shape_id": 999999}, "x")
    assert ok is False


def test_set_alt_text_returns_false_for_non_shape_target(tmp_path):
    from pptx_a11y.appliers import set_alt_text

    prs = _prs_with_picture(tmp_path)
    # run-level target: picture has no text frame, resolve returns None
    ok = set_alt_text(prs, {"slide": 0, "shape_id": 999, "para": 0, "run": 0}, "x")
    assert ok is False


# ---------------------------------------------------------------------------
# mark_decorative
# ---------------------------------------------------------------------------

def test_mark_decorative_adds_adec_element(tmp_path):
    from pptx_a11y.appliers import mark_decorative
    from lxml import etree

    prs = _prs_with_picture(tmp_path)
    pic, target = _pic_and_target(prs)
    ok = mark_decorative(prs, target, None)
    assert ok is True
    cNvPr = pic._element._nvXxPr.cNvPr
    ns = "http://schemas.microsoft.com/office/drawing/2017/decorative"
    found = cNvPr.find(f"{{{ns}}}decorative")
    assert found is not None
    assert found.get("val") == "1"


def test_mark_decorative_returns_false_for_unknown_shape(tmp_path):
    from pptx_a11y.appliers import mark_decorative

    prs = _prs_with_picture(tmp_path)
    ok = mark_decorative(prs, {"slide": 0, "shape_id": 999999}, None)
    assert ok is False


# ---------------------------------------------------------------------------
# set_title
# ---------------------------------------------------------------------------

def test_set_title_sets_slide_title_text(tmp_path):
    from pptx_a11y.appliers import set_title

    prs = _prs_with_picture(tmp_path)
    target = {"slide": 0, "scope": "slide_title"}
    ok = set_title(prs, target, "New Title")
    assert ok is True
    assert prs.slides[0].shapes.title.text == "New Title"


def test_set_title_returns_false_when_no_title_placeholder(tmp_path):
    from pptx_a11y.appliers import set_title

    prs = Presentation()
    # layout 6 = blank — no title placeholder
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    assert slide.shapes.title is None
    path = str(tmp_path / "notitle.pptx")
    prs.save(path)
    prs2 = Presentation(path)
    ok = set_title(prs2, {"slide": 0, "scope": "slide_title"}, "X")
    assert ok is False


# ---------------------------------------------------------------------------
# set_doc_title
# ---------------------------------------------------------------------------

def test_set_doc_title_sets_core_properties(tmp_path):
    from pptx_a11y.appliers import set_doc_title

    prs = _prs_with_picture(tmp_path)
    ok = set_doc_title(prs, {"scope": "document"}, "My Presentation")
    assert ok is True
    assert prs.core_properties.title == "My Presentation"


# ---------------------------------------------------------------------------
# set_doc_language
# ---------------------------------------------------------------------------

def test_set_doc_language_sets_core_properties(tmp_path):
    from pptx_a11y.appliers import set_doc_language

    prs = _prs_with_picture(tmp_path)
    ok = set_doc_language(prs, {"scope": "document"}, "fr-FR")
    assert ok is True
    assert prs.core_properties.language == "fr-FR"


# ---------------------------------------------------------------------------
# set_link_text (hyperlink address must survive)
# ---------------------------------------------------------------------------

def test_set_link_text_updates_run_text(tmp_path):
    from pptx_a11y.appliers import set_link_text

    prs = _prs_with_run(tmp_path)
    run, target = _run_and_target(prs)
    original_address = run.hyperlink.address
    assert original_address  # sanity: deck_with_issues wires a hyperlink
    ok = set_link_text(prs, target, "Visit example website")
    assert ok is True
    assert run.text == "Visit example website"


def test_set_link_text_preserves_hyperlink_address(tmp_path):
    from pptx_a11y.appliers import set_link_text

    prs = _prs_with_run(tmp_path)
    run, target = _run_and_target(prs)
    original_address = run.hyperlink.address
    set_link_text(prs, target, "Descriptive link text")
    assert run.hyperlink.address == original_address


def test_set_link_text_returns_false_for_bad_target(tmp_path):
    from pptx_a11y.appliers import set_link_text

    prs = _prs_with_run(tmp_path)
    ok = set_link_text(prs, {"slide": 9, "shape_id": 1, "para": 0, "run": 0}, "x")
    assert ok is False


# ---------------------------------------------------------------------------
# set_table_header
# ---------------------------------------------------------------------------

def test_set_table_header_sets_first_row_true(tmp_path):
    from pptx_a11y.appliers import set_table_header

    prs = _prs_with_table(tmp_path)
    shape, target = _table_shape_target(prs)
    ok = set_table_header(prs, target, None)
    assert ok is True
    assert shape.table.first_row is True


def test_set_table_header_returns_false_for_non_table_shape(tmp_path):
    from pptx_a11y.appliers import set_table_header

    prs = _prs_with_picture(tmp_path)
    pic, target = _pic_and_target(prs)
    ok = set_table_header(prs, target, None)
    assert ok is False


def test_set_table_header_returns_false_for_unknown_shape(tmp_path):
    from pptx_a11y.appliers import set_table_header

    prs = _prs_with_table(tmp_path)
    ok = set_table_header(prs, {"slide": 0, "shape_id": 999999}, None)
    assert ok is False


# ---------------------------------------------------------------------------
# apply_contrast_color
# ---------------------------------------------------------------------------

def test_apply_contrast_color_sets_rgb_from_list(tmp_path):
    from pptx_a11y.appliers import apply_contrast_color
    from pptx.dml.color import RGBColor

    prs = _prs_with_run(tmp_path)
    run, target = _run_and_target(prs)
    ok = apply_contrast_color(prs, target, [0, 0, 0])
    assert ok is True
    assert run.font.color.rgb == RGBColor(0, 0, 0)


def test_apply_contrast_color_accepts_hex_string(tmp_path):
    from pptx_a11y.appliers import apply_contrast_color
    from pptx.dml.color import RGBColor

    prs = _prs_with_run(tmp_path)
    run, target = _run_and_target(prs)
    ok = apply_contrast_color(prs, target, "#1a2b3c")
    assert ok is True
    assert run.font.color.rgb == RGBColor(0x1a, 0x2b, 0x3c)


def test_apply_contrast_color_returns_false_for_bad_target(tmp_path):
    from pptx_a11y.appliers import apply_contrast_color

    prs = _prs_with_run(tmp_path)
    ok = apply_contrast_color(prs, {"slide": 99, "shape_id": 1, "para": 0, "run": 0}, [0, 0, 0])
    assert ok is False


# ---------------------------------------------------------------------------
# bump_font_size
# ---------------------------------------------------------------------------

def test_bump_font_size_sets_pt_from_value(tmp_path):
    from pptx_a11y.appliers import bump_font_size

    prs = _prs_with_run(tmp_path)
    run, target = _run_and_target(prs)
    ok = bump_font_size(prs, target, 24)
    assert ok is True
    assert run.font.size.pt == 24


def test_bump_font_size_enforces_minimum_18(tmp_path):
    from pptx_a11y.appliers import bump_font_size

    prs = _prs_with_run(tmp_path)
    run, target = _run_and_target(prs)
    ok = bump_font_size(prs, target, 10)  # below minimum
    assert ok is True
    assert run.font.size.pt >= 18


def test_bump_font_size_returns_false_for_bad_target(tmp_path):
    from pptx_a11y.appliers import bump_font_size

    prs = _prs_with_run(tmp_path)
    ok = bump_font_size(prs, {"slide": 99, "shape_id": 1, "para": 0, "run": 0}, 24)
    assert ok is False


# ---------------------------------------------------------------------------
# APPLIERS registry
# ---------------------------------------------------------------------------

def test_appliers_registry_has_all_expected_keys():
    from pptx_a11y.appliers import APPLIERS

    expected = {
        "set_alt_text",
        "mark_decorative",
        "set_title",
        "set_doc_title",
        "set_doc_language",
        "set_link_text",
        "set_table_header",
        "apply_contrast_color",
        "bump_font_size",
    }
    assert expected.issubset(set(APPLIERS.keys()))


def test_appliers_registry_values_are_callable():
    from pptx_a11y.appliers import APPLIERS

    for name, fn in APPLIERS.items():
        assert callable(fn), f"APPLIERS[{name!r}] is not callable"


# ---------------------------------------------------------------------------
# apply_plan
# ---------------------------------------------------------------------------

def test_apply_plan_returns_ok_for_valid_plan(tmp_path):
    from pptx_a11y.appliers import apply_plan

    prs = _prs_with_run(tmp_path)
    _, target = _run_and_target(prs)
    plan = [{"action": "bump_font_size", "target": target, "value": 24}]
    results = apply_plan(prs, plan)
    assert len(results) == 1
    assert results[0]["action"] == "bump_font_size"
    assert results[0]["ok"] is True


def test_apply_plan_skips_unresolvable_target_without_raising(tmp_path):
    from pptx_a11y.appliers import apply_plan

    prs = _prs_with_run(tmp_path)
    bad_target = {"slide": 99, "shape_id": 999}
    plan = [{"action": "set_alt_text", "target": bad_target, "value": "x"}]
    results = apply_plan(prs, plan)
    assert len(results) == 1
    assert results[0]["ok"] is False


def test_apply_plan_unknown_action_records_false(tmp_path):
    from pptx_a11y.appliers import apply_plan

    prs = _prs_with_picture(tmp_path)
    plan = [{"action": "nonexistent_action", "target": {}, "value": None}]
    results = apply_plan(prs, plan)
    assert results[0]["ok"] is False


def test_apply_plan_one_bad_item_doesnt_abort(tmp_path):
    from pptx_a11y.appliers import apply_plan

    prs = _prs_with_run(tmp_path)
    _, run_target = _run_and_target(prs)
    plan = [
        {"action": "nonexistent_action", "target": {}, "value": None},  # bad
        {"action": "bump_font_size", "target": run_target, "value": 20},  # good
    ]
    results = apply_plan(prs, plan)
    assert len(results) == 2
    assert results[0]["ok"] is False
    assert results[1]["ok"] is True


def test_apply_plan_handles_empty_plan(tmp_path):
    from pptx_a11y.appliers import apply_plan

    prs = _prs_with_picture(tmp_path)
    results = apply_plan(prs, [])
    assert results == []


def test_apply_plan_multiple_actions(tmp_path):
    from pptx_a11y.appliers import apply_plan

    prs = _prs_with_run(tmp_path)
    pic_shape = next(s for s in prs.slides[0].shapes if s.shape_type == 13)
    pic_target = {"slide": 0, "shape_id": pic_shape.shape_id}
    _, run_target = _run_and_target(prs)

    plan = [
        {"action": "set_alt_text", "target": pic_target, "value": "A diagram"},
        {"action": "set_doc_title", "target": {"scope": "document"}, "value": "My Deck"},
        {"action": "bump_font_size", "target": run_target, "value": 20},
    ]
    results = apply_plan(prs, plan)
    assert all(r["ok"] for r in results)
    assert prs.core_properties.title == "My Deck"
    assert pic_shape._element._nvXxPr.cNvPr.get("descr") == "A diagram"
