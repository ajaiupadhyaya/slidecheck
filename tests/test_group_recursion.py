"""Group-nested shapes must be reached by the text/table checks, not just the
top-level ones. iter_runs and the tables check now descend into groups."""
from pptx import Presentation
from pptx.util import Inches, Pt

from pptx_a11y.checks.font_size import check as font_check
from pptx_a11y.checks.link_text import check as link_check
from pptx_a11y.textutil import iter_runs


def _deck_with_grouped_textbox(path, text, size_pt=None, link=None):
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    group = s.shapes.add_group_shape()
    tb = group.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = text
    if size_pt is not None:
        run.font.size = Pt(size_pt)
    if link is not None:
        run.hyperlink.address = link
    prs.save(path)
    return path


def test_iter_runs_descends_into_groups(tmp_path):
    p = _deck_with_grouped_textbox(str(tmp_path / "g.pptx"), "inside a group")
    texts = [run.text for _i, _shape, _para, run in iter_runs(Presentation(p))]
    assert "inside a group" in texts


def test_font_size_check_sees_grouped_text(tmp_path):
    p = _deck_with_grouped_textbox(str(tmp_path / "g.pptx"), "tiny grouped text", size_pt=10)
    assert any(f.check_id == "font_size" for f in font_check(Presentation(p)))


def test_link_text_check_sees_grouped_link(tmp_path):
    p = _deck_with_grouped_textbox(
        str(tmp_path / "g.pptx"), "click here", link="https://example.com"
    )
    assert any(f.check_id == "link_text" for f in link_check(Presentation(p)))
