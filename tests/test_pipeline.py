import os
from pptx_a11y.pipeline import process_file, unique_path
from tests.fixtures.build import deck_with_issues


class _StubDescriber:
    def describe(self, image_bytes, media_type, context):
        return "A described image."

    def suggest_text(self, prompt):
        return "Stub Title"


def test_process_file_produces_outputs_and_marks_fixed(tmp_path):
    src = deck_with_issues(str(tmp_path / "deck.pptx"))
    result = process_file(src, _StubDescriber())
    assert result.error is None
    assert os.path.exists(result.output_path)
    assert result.output_path.endswith("_accessible.pptx")
    assert os.path.exists(str(tmp_path / "deck_a11y_report.html"))
    assert os.path.exists(str(tmp_path / "deck_a11y_report.json"))
    # original untouched: still no alt text in the source file
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    src_prs = Presentation(src)
    pics = [sh for sl in src_prs.slides for sh in sl.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert not (pics[0]._element._nvXxPr.cNvPr.get("descr") or "")
    # an alt_text finding got marked auto_fixed
    assert any(f.check_id == "alt_text" and f.auto_fixed for f in result.findings)


def test_corrupt_file_returns_error_result(tmp_path):
    bad = tmp_path / "bad.pptx"
    bad.write_bytes(b"nope")
    result = process_file(str(bad), _StubDescriber())
    assert result.error is not None
    assert result.output_path is None


def test_unique_path_disambiguates(tmp_path):
    p = tmp_path / "deck_accessible.pptx"
    p.write_bytes(b"x")
    assert unique_path(str(p)).endswith("deck_accessible_1.pptx")
