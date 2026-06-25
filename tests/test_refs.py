"""Tests for refs.py: shape_target, run_target, resolve_target."""
import base64

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from pptx_a11y.refs import resolve_target, run_target, shape_target
from tests.fixtures.build import deck_with_issues

# ---------------------------------------------------------------------------
# shape_target / run_target helpers
# ---------------------------------------------------------------------------

def test_shape_target_returns_dict_with_slide_and_shape_id(tmp_path):
    prs = Presentation(deck_with_issues(str(tmp_path / "d.pptx")))
    slide0 = prs.slides[0]
    pic = next(s for s in slide0.shapes if s.shape_type == 13)  # PICTURE
    t = shape_target(0, pic)
    assert t["slide"] == 0
    assert t["shape_id"] == pic.shape_id


def test_run_target_returns_dict_with_all_keys(tmp_path):
    prs = Presentation(deck_with_issues(str(tmp_path / "d.pptx")))
    slide1 = prs.slides[1]
    shape = next(s for s in slide1.shapes if s.has_text_frame)
    t = run_target(1, shape, 0, 0)
    assert t == {"slide": 1, "shape_id": shape.shape_id, "para": 0, "run": 0}


# ---------------------------------------------------------------------------
# resolve_target: shape round-trip
# ---------------------------------------------------------------------------

def test_resolve_shape_target_roundtrip(tmp_path):
    prs = Presentation(deck_with_issues(str(tmp_path / "d.pptx")))
    slide0 = prs.slides[0]
    pic = next(s for s in slide0.shapes if s.shape_type == 13)  # PICTURE
    t = shape_target(0, pic)
    got = resolve_target(prs, t)
    assert got is not None and got.shape_id == pic.shape_id


def test_resolve_run_target_roundtrip(tmp_path):
    prs = Presentation(deck_with_issues(str(tmp_path / "d.pptx")))
    slide1 = prs.slides[1]
    # find a shape with at least one run on the first para
    shape = next(
        s for s in slide1.shapes
        if s.has_text_frame and s.text_frame.paragraphs
        and s.text_frame.paragraphs[0].runs
    )
    run = shape.text_frame.paragraphs[0].runs[0]
    t = run_target(1, shape, 0, 0)
    got = resolve_target(prs, t)
    assert got is not None
    assert got.text == run.text


# ---------------------------------------------------------------------------
# resolve_target: document scope
# ---------------------------------------------------------------------------

def test_resolve_document_scope(tmp_path):
    prs = Presentation(deck_with_issues(str(tmp_path / "d.pptx")))
    assert resolve_target(prs, {"scope": "document"}) is prs


def test_resolve_document_scope_with_field(tmp_path):
    prs = Presentation(deck_with_issues(str(tmp_path / "d.pptx")))
    assert resolve_target(prs, {"scope": "document", "field": "title"}) is prs


def test_resolve_slide_title_scope(tmp_path):
    """slide_title scope (has no shape_id) resolves to prs (document-level)."""
    prs = Presentation(deck_with_issues(str(tmp_path / "d.pptx")))
    t = {"slide": 0, "scope": "slide_title"}
    result = resolve_target(prs, t)
    assert result is prs


# ---------------------------------------------------------------------------
# resolve_target: unknown / out-of-range returns None
# ---------------------------------------------------------------------------

def test_resolve_unknown_returns_none(tmp_path):
    prs = Presentation(deck_with_issues(str(tmp_path / "d.pptx")))
    assert resolve_target(prs, {"slide": 0, "shape_id": 999999}) is None


def test_resolve_bad_slide_index_returns_none(tmp_path):
    prs = Presentation(deck_with_issues(str(tmp_path / "d.pptx")))
    assert resolve_target(prs, {"slide": 999, "shape_id": 1}) is None


def test_resolve_bad_run_index_returns_none(tmp_path):
    prs = Presentation(deck_with_issues(str(tmp_path / "d.pptx")))
    slide0 = prs.slides[0]
    pic = next(s for s in slide0.shapes if s.shape_type == 13)  # PICTURE (no text frame)
    t = {"slide": 0, "shape_id": pic.shape_id, "para": 0, "run": 0}
    assert resolve_target(prs, t) is None


def test_resolve_empty_target_returns_none(tmp_path):
    prs = Presentation(deck_with_issues(str(tmp_path / "d.pptx")))
    assert resolve_target(prs, {}) is None


# ---------------------------------------------------------------------------
# resolve_target: group descent
# ---------------------------------------------------------------------------

def test_resolve_shape_in_group(tmp_path):
    """A shape nested inside a group can be resolved by shape_id."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb1 = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    tb2 = slide.shapes.add_textbox(Inches(3), Inches(1), Inches(2), Inches(1))
    grp = slide.shapes.add_group_shape([tb1, tb2])
    inner_shape = grp.shapes[0]
    inner_id = inner_shape.shape_id
    path = str(tmp_path / "grp.pptx")
    prs.save(path)

    prs2 = Presentation(path)
    t = {"slide": 0, "shape_id": inner_id}
    got = resolve_target(prs2, t)
    assert got is not None
    assert got.shape_id == inner_id
