"""Tests for pptx_a11y.web.export_service — TDD RED phase first."""
from __future__ import annotations

import io
import tempfile

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from pptx_a11y.alt_text_ai import NullDescriber
from pptx_a11y.web.analyze_service import analyze_upload
from pptx_a11y.web.export_service import export_with_plan
from tests.fixtures.build import deck_with_issues, clean_deck


def _bytes(tmp_path, builder, name="src.pptx"):
    p = builder(str(tmp_path / name))
    with open(p, "rb") as fh:
        return fh.read()


def _pic_target(tmp_path):
    """Return (data bytes, picture target dict) for the picture on slide 0 of deck_with_issues."""
    data = _bytes(tmp_path, deck_with_issues)
    prs = Presentation(io.BytesIO(data))
    slide0 = prs.slides[0]
    pic = next(s for s in slide0.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
    target = {"slide": 0, "shape_id": pic.shape_id}
    return data, target


# ---------------------------------------------------------------------------
# Happy path: set_alt_text on the picture
# ---------------------------------------------------------------------------

def test_export_with_plan_returns_fixed_bytes(tmp_path):
    data, target = _pic_target(tmp_path)
    plan = [{"action": "set_alt_text", "target": target, "value": "A red square"}]
    result = export_with_plan("lecture.pptx", data, plan)
    assert result["error"] is None
    assert result["fixed_bytes"] is not None


def test_export_with_plan_fixed_bytes_is_valid_pptx(tmp_path):
    data, target = _pic_target(tmp_path)
    plan = [{"action": "set_alt_text", "target": target, "value": "A red square"}]
    result = export_with_plan("lecture.pptx", data, plan)
    Presentation(io.BytesIO(result["fixed_bytes"]))  # must not raise


def test_export_with_plan_alt_text_applied_in_fixed_pptx(tmp_path):
    data, target = _pic_target(tmp_path)
    plan = [{"action": "set_alt_text", "target": target, "value": "A red square"}]
    result = export_with_plan("lecture.pptx", data, plan)
    prs2 = Presentation(io.BytesIO(result["fixed_bytes"]))
    slide0 = prs2.slides[0]
    pic = next(s for s in slide0.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
    assert pic._element._nvXxPr.cNvPr.get("descr") == "A red square"


def test_export_with_plan_applied_shows_ok_true(tmp_path):
    data, target = _pic_target(tmp_path)
    plan = [{"action": "set_alt_text", "target": target, "value": "A red square"}]
    result = export_with_plan("lecture.pptx", data, plan)
    assert result["applied"] == [{"action": "set_alt_text", "ok": True}]


def test_export_with_plan_report_html_is_non_empty(tmp_path):
    data, target = _pic_target(tmp_path)
    plan = [{"action": "set_alt_text", "target": target, "value": "A red square"}]
    result = export_with_plan("lecture.pptx", data, plan)
    assert result["report_html"]
    assert "<" in result["report_html"]


def test_export_with_plan_fixed_filename_has_accessible_suffix(tmp_path):
    data, target = _pic_target(tmp_path)
    plan = [{"action": "set_alt_text", "target": target, "value": "A red square"}]
    result = export_with_plan("lecture.pptx", data, plan)
    assert result["fixed_filename"] == "lecture_accessible.pptx"


def test_export_with_plan_untargeted_element_is_unchanged(tmp_path):
    """The 'click here' run on slide 1 must not be mutated by a plan that only
    touches the picture on slide 0."""
    data, target = _pic_target(tmp_path)
    plan = [{"action": "set_alt_text", "target": target, "value": "A red square"}]
    result = export_with_plan("lecture.pptx", data, plan)
    prs2 = Presentation(io.BytesIO(result["fixed_bytes"]))
    slide1 = prs2.slides[1]
    # Find the hyperlink run on slide 1
    link_run = None
    for shape in slide1.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.hyperlink.address:
                        link_run = run
    assert link_run is not None, "hyperlink run must still exist on slide 1"
    assert link_run.text == "click here", \
        f"untargeted run text must be unchanged, got: {link_run.text!r}"


def test_export_with_plan_original_data_not_mutated(tmp_path):
    """Original bytes must represent the pre-fix deck (picture descr still empty)."""
    data, target = _pic_target(tmp_path)
    plan = [{"action": "set_alt_text", "target": target, "value": "A red square"}]
    export_with_plan("lecture.pptx", data, plan)
    # Re-open original
    prs_orig = Presentation(io.BytesIO(data))
    slide0 = prs_orig.slides[0]
    pic = next(s for s in slide0.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
    # Original had no alt text
    assert pic._element._nvXxPr.cNvPr.get("descr") == "", \
        "original data must not be mutated by export_with_plan"


def test_export_with_plan_returns_filename(tmp_path):
    data, target = _pic_target(tmp_path)
    plan = [{"action": "set_alt_text", "target": target, "value": "A red square"}]
    result = export_with_plan("lecture.pptx", data, plan)
    assert result["filename"] == "lecture.pptx"


def test_export_with_plan_report_html_uses_original_filename(tmp_path):
    data, target = _pic_target(tmp_path)
    plan = [{"action": "set_alt_text", "target": target, "value": "A red square"}]
    result = export_with_plan("My Lecture.pptx", data, plan)
    assert "My Lecture" in result["report_html"]
    assert tempfile.gettempdir() not in result["report_html"]


def test_export_with_plan_empty_plan_returns_unchanged_deck(tmp_path):
    data = _bytes(tmp_path, clean_deck)
    result = export_with_plan("clean.pptx", data, [])
    assert result["error"] is None
    assert result["applied"] == []
    Presentation(io.BytesIO(result["fixed_bytes"]))  # still valid pptx


def test_export_with_plan_leaves_no_files(tmp_path, monkeypatch):
    """Nothing persists after export_with_plan returns."""
    import os

    monkeypatch.chdir(tmp_path)
    data = _bytes(tmp_path, clean_deck, "keep.pptx")
    export_with_plan("keep.pptx", data, [])
    leftover = [n for n in os.listdir(tmp_path) if n.endswith((".pptx", ".html")) and n != "keep.pptx"]
    assert leftover == []


# ---------------------------------------------------------------------------
# Error path: corrupt bytes
# ---------------------------------------------------------------------------

def test_export_with_plan_corrupt_bytes_returns_error(tmp_path):
    result = export_with_plan("broken.pptx", b"not a pptx", [])
    assert result["filename"] == "broken.pptx"
    assert result["error"] is not None
    assert result["fixed_bytes"] is None
    assert result["fixed_filename"] is None
    assert result["report_html"] is None
    assert result["applied"] == []


def test_export_with_plan_corrupt_bytes_hides_temp_path(tmp_path):
    result = export_with_plan("my file.pptx", b"not a pptx", [])
    err = result["error"]
    assert err
    assert tempfile.gettempdir() not in err


# ---------------------------------------------------------------------------
# Roundtrip via analyze → export
# ---------------------------------------------------------------------------

def test_analyze_then_export_roundtrip(tmp_path):
    """Analyze a deck, build a plan from its findings, export, re-open, verify."""
    data = _bytes(tmp_path, deck_with_issues)
    analysis_result = analyze_upload("lecture.pptx", data, NullDescriber())
    findings = analysis_result["analysis"]["findings"]
    # Build plan for any fixable finding
    plan = [
        {"action": f["fix_action"], "target": f["target"], "value": f.get("suggested_value") or "Fixed value"}
        for f in findings
        if f.get("fixable") and f.get("fix_action") in ("set_alt_text",)
    ]
    assert plan, "expected at least one set_alt_text finding to fix"
    result = export_with_plan("lecture.pptx", data, plan)
    assert result["error"] is None
    Presentation(io.BytesIO(result["fixed_bytes"]))  # valid pptx
    assert any(a["ok"] for a in result["applied"])
